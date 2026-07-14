from __future__ import annotations

import argparse
import statistics
from pathlib import Path
from datetime import datetime

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from scripts.export_figures import generate_summary_figures
from utils.aligned_utils import aligned_window_tag, compute_aligned_window, resolve_summary_windows
from utils.analysis_mode_utils import resolve_effective_analysis_mode
from utils.event_utils import read_light_intervals
from utils.file_id_utils import canonicalize_file_id, infer_has_light_from_identifiers, legacy_file_id_candidates
from utils.logging_utils import PipelineLogger
from utils.path_utils import load_yaml, resolve_project_paths
from utils.table_utils import normalize_include_column, normalize_stim_schedule, read_table


def _aligned_cfg(config: dict) -> dict:
    return config.get("aligned_rate", config.get("neuroexplorer", {}).get("aligned_rate", {}))


def _add_textbox(slide, left: float, top: float, width: float, height: float, text: str, font_size: int = 16) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(font_size)


def _clean_meta_value(value, *, max_len: int | None = None) -> str:
    if value is None:
        return "-"
    text = str(value).strip()
    if not text or text.lower() == "nan":
        return "-"
    if max_len is not None and len(text) > max_len:
        return text[: max_len - 3].rstrip() + "..."
    return text


def _add_metadata_two_columns(slide, left: float, top: float, width: float, height: float, lines: tuple[list[str], list[str]], font_size: int = 10) -> None:
    column_gap = 0.2
    column_width = (width - column_gap) / 2
    _add_textbox(slide, left, top, column_width, height, "\n".join(lines[0]), font_size=font_size)
    _add_textbox(slide, left + column_width + column_gap, top, column_width, height, "\n".join(lines[1]), font_size=font_size)


def _add_picture_contain(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    with Image.open(image_path) as image:
        img_width_px, img_height_px = image.size

    scale = min(width / img_width_px, height / img_height_px)
    render_width = img_width_px * scale
    render_height = img_height_px * scale
    render_left = left + (width - render_width) / 2
    render_top = top + (height - render_height) / 2
    slide.shapes.add_picture(
        str(image_path),
        Inches(render_left),
        Inches(render_top),
        width=Inches(render_width),
        height=Inches(render_height),
    )


def _window_summary(config: dict, stim_sub) -> dict:
    aligned_cfg = _aligned_cfg(config)
    if stim_sub.empty or not bool(stim_sub["has_light_bool"].any()):
        return {
            "window_mode": aligned_cfg.get("window_mode", "configured_windows"),
            "summary_window_mode": "-",
            "absolute_window": "-",
            "aligned_window": "not applicable",
            "baseline_window": "not applicable",
            "light_window": "none",
            "post_window": "not applicable",
        }
    windows = [
        compute_aligned_window(
            light_on_s=float(row.light_on_s),
            light_off_s=float(row.light_off_s),
            duration_s=float(row.duration_s),
            aligned_cfg=aligned_cfg,
        )
        for row in stim_sub.itertuples(index=False)
    ]
    if not windows:
        return {
            "window_mode": aligned_cfg.get("window_mode", "configured_windows"),
            "summary_window_mode": "-",
            "absolute_window": "-",
            "aligned_window": "-",
            "baseline_window": "-",
            "light_window": "-",
            "post_window": "-",
        }
    durations = [window["duration_s"] for window in windows]
    unique_duration = {round(value, 6) for value in durations}
    if len(unique_duration) == 1:
        duration_value = durations[0]
        sample = windows[0]
        summary_windows = resolve_summary_windows(duration_value, aligned_cfg)
        return {
            "window_mode": sample["window_mode"],
            "summary_window_mode": summary_windows["summary_window_mode"],
            "absolute_window": f"{sample['abs_start_s']:g} to {sample['abs_end_s']:g} s",
            "aligned_window": f"{sample['aligned_x_min_s']:g} to {sample['aligned_x_max_s']:g} s",
            "baseline_window": f"{summary_windows['baseline'][0]:g} to {summary_windows['baseline'][1]:g} s",
            "light_window": f"{summary_windows['light'][0]:g} to {summary_windows['light'][1]:g} s",
            "post_window": f"{summary_windows['post'][0]:g} to {summary_windows['post'][1]:g} s",
        }
    min_aligned = min(window["aligned_x_min_s"] for window in windows)
    max_aligned = max(window["aligned_x_max_s"] for window in windows)
    min_abs = min(window["abs_start_s"] for window in windows)
    max_abs = max(window["abs_end_s"] for window in windows)
    median_duration = statistics.median(durations)
    summary_windows = resolve_summary_windows(float(median_duration), aligned_cfg)
    return {
        "window_mode": aligned_cfg.get("window_mode", "configured_windows"),
        "summary_window_mode": summary_windows["summary_window_mode"],
        "absolute_window": f"{min_abs:g} to {max_abs:g} s",
        "aligned_window": f"{min_aligned:g} to {max_aligned:g} s (duration range {min(durations):g}-{max(durations):g} s)",
        "baseline_window": f"{summary_windows['baseline'][0]:g} to {summary_windows['baseline'][1]:g} s (median duration basis)",
        "light_window": f"{summary_windows['light'][0]:g} to {summary_windows['light'][1]:g} s (median duration basis)",
        "post_window": f"{summary_windows['post'][0]:g} to {summary_windows['post'][1]:g} s (median duration basis)",
    }


def _metadata_lines(
    config: dict,
    file_id: str,
    unit_id: str,
    original_name: str,
    unit_note: str,
    n_events: int,
    duration_summary: str,
    light_on_summary: str,
    light_off_summary: str,
    window_summary: dict,
    *,
    has_light: bool,
    event_group: str,
    unit_source: str,
    effective_mode: str,
) -> tuple[list[str], list[str]]:
    aligned_cfg = _aligned_cfg(config)
    full_cfg = config["neuroexplorer"]["fullrate"]
    analysis_mode = config.get("analysis", {}).get("mode", "fullrate_aligned")
    qc_note = _clean_meta_value(unit_note, max_len=120)
    left_lines = [
        f"file_id: {_clean_meta_value(file_id)}",
        f"unit_id: {_clean_meta_value(unit_id)}",
        f"original_name: {_clean_meta_value(original_name)}",
        f"unit_source: {_clean_meta_value(unit_source)}",
        f"analysis_mode: {_clean_meta_value(analysis_mode)}",
        f"effective_mode: {_clean_meta_value(effective_mode)}",
        f"has_light: {'yes' if has_light else 'no'}",
        f"event_group: {_clean_meta_value(event_group or ('nolight' if not has_light else '-'))}",
        f"Full-session bin width: {_clean_meta_value(full_cfg.get('bin_width_s'))} s",
        f"QC note: {qc_note}",
    ]
    if not has_light:
        return left_lines, [
            "alignment: not applicable",
            "light_band: none",
            "NeuroExplorer event required: no",
            "PPTX page: no-light control",
        ]
    event_required = "no" if effective_mode == "fullrate_aligned" else "yes"
    right_lines = [
        f"window_mode: {_clean_meta_value(window_summary.get('window_mode'))}",
        f"summary_window_mode: {_clean_meta_value(window_summary.get('summary_window_mode'))}",
        f"absolute_window: {_clean_meta_value(window_summary.get('absolute_window'))}",
        f"aligned_window: {_clean_meta_value(window_summary.get('aligned_window'))}",
        f"baseline_window: {_clean_meta_value(window_summary.get('baseline_window'))}",
        f"light_window: {_clean_meta_value(window_summary.get('light_window'))}",
        f"post_window: {_clean_meta_value(window_summary.get('post_window'))}",
        f"light_on_times: {_clean_meta_value(light_on_summary)}",
        f"light_off_times: {_clean_meta_value(light_off_summary)}",
        f"duration_s: {_clean_meta_value(duration_summary)}",
        "alignment_source: stim_schedule_master",
        f"NeuroExplorer event required: {event_required}",
    ]
    if effective_mode == "fullrate_aligned":
        right_lines.append("No NeuroExplorer event variable required; alignment source = stim_schedule_master.xlsx.")
    return left_lines, right_lines


def _interval_metadata(config: dict, paths: dict, file_id: str, stim_sub) -> tuple[str, str, str, int]:
    if stim_sub.empty or not bool(stim_sub["has_light_bool"].any()):
        return "-", "-", "-", 0
    interval_cfg = config.get("neuroexplorer", {}).get("interval", {})
    interval_path = paths["events_export_dir"] / interval_cfg.get("interval_csv_pattern", "{file_id}_Light_Interval.csv").format(file_id=file_id)
    if interval_path.exists():
        intervals = read_light_intervals(interval_path, delimiter=interval_cfg.get("delimiter", ","))
    else:
        intervals = [(float(row.light_on_s), float(row.light_off_s)) for row in stim_sub.itertuples(index=False)]
    light_on = [start for start, _ in intervals]
    light_off = [end for _, end in intervals]
    durations = [end - start for start, end in intervals]
    if len({round(value, 6) for value in durations}) == 1:
        duration_summary = str(durations[0]).rstrip("0").rstrip(".") if not float(durations[0]).is_integer() else str(int(durations[0]))
    else:
        duration_summary = f"{min(durations)}-{max(durations)}"
    light_on_summary = ", ".join(str(int(x)) if float(x).is_integer() else str(x) for x in light_on)
    light_off_summary = ", ".join(str(int(x)) if float(x).is_integer() else str(x) for x in light_off)
    return duration_summary, light_on_summary, light_off_summary, len(intervals)


def _expected_fullrate_path(config: dict, paths: dict, file_id: str) -> Path:
    return paths["nex_fullrate_dir"] / config["neuroexplorer"]["export"]["expected_fullrate_pattern"].format(
        file_id=file_id,
        bin_width_s=config["neuroexplorer"]["fullrate"]["bin_width_s"],
    )


def _figure_path_with_legacy(
    *,
    paths: dict,
    config: dict,
    file_id: str,
    pl2_file: str,
    unit_id: str,
    figure_dir_key: str,
    suffix: str,
    logger: PipelineLogger,
) -> Path:
    canonical_path = paths[figure_dir_key] / f"{file_id}_{unit_id}_{suffix}"
    if canonical_path.exists():
        return canonical_path
    for candidate_file_id in legacy_file_id_candidates(file_id, pl2_file, config)[1:]:
        legacy_path = paths[figure_dir_key] / f"{candidate_file_id}_{unit_id}_{suffix}"
        if legacy_path.exists():
            logger.log(
                "build_pptx",
                file_id,
                str(legacy_path),
                str(canonical_path),
                "warning",
                "Using legacy file_id figure path; please rerun export_figures after canonicalizing file_id.",
            )
            return legacy_path
    return canonical_path


def _units_from_fullrate_csv(config: dict, paths: dict, file_id: str) -> list[dict]:
    fullrate_csv = next(
        (
            paths["nex_fullrate_dir"] / config["neuroexplorer"]["export"]["expected_fullrate_pattern"].format(
                file_id=candidate,
                bin_width_s=config["neuroexplorer"]["fullrate"]["bin_width_s"],
            )
            for candidate in legacy_file_id_candidates(file_id, None, config)
            if (
                paths["nex_fullrate_dir"] / config["neuroexplorer"]["export"]["expected_fullrate_pattern"].format(
                    file_id=candidate,
                    bin_width_s=config["neuroexplorer"]["fullrate"]["bin_width_s"],
                )
            ).exists()
        ),
        None,
    )
    if fullrate_csv is None:
        return []
    try:
        fullrate_df = read_table(fullrate_csv)
    except Exception:
        return []
    if "unit_id" not in fullrate_df.columns:
        return []
    rows = []
    original_name_by_unit = {}
    if "original_name" in fullrate_df.columns:
        original_name_by_unit = {
            str(row.unit_id): str(row.original_name)
            for row in fullrate_df[["unit_id", "original_name"]].dropna(subset=["unit_id"]).drop_duplicates().itertuples(index=False)
        }
    for unit_id in sorted({str(value).strip() for value in fullrate_df["unit_id"].dropna().astype(str) if str(value).strip()}):
        rows.append(
            {
                config["project"]["file_id_column"]: file_id,
                "unit_id": unit_id,
                "original_name": original_name_by_unit.get(unit_id, unit_id),
                "include": "yes",
                "include_bool": True,
                "note": "",
                "unit_source": "fullrate_csv_fallback",
            }
        )
    return rows


def _augment_units_from_fullrate_fallback(config: dict, paths: dict, stim_df: pd.DataFrame, included_df: pd.DataFrame, logger: PipelineLogger) -> pd.DataFrame:
    file_id_column = config["project"]["file_id_column"]
    existing_file_ids = set(included_df[file_id_column].astype(str).tolist()) if not included_df.empty else set()
    fallback_rows: list[dict] = []
    for file_id in stim_df[file_id_column].astype(str).drop_duplicates().tolist():
        if file_id in existing_file_ids:
            continue
        rows = _units_from_fullrate_csv(config, paths, file_id)
        if not rows:
            continue
        fallback_rows.extend(rows)
        logger.log(
            "build_pptx",
            file_id,
            str(_expected_fullrate_path(config, paths, file_id)),
            "",
            "warning",
            "unit_quality_table missing rows; using fullrate CSV unit_id fallback for PPTX.",
        )
    if not fallback_rows:
        if "unit_source" not in included_df.columns:
            included_df = included_df.copy()
            included_df["unit_source"] = "unit_quality_table"
        return included_df
    base_df = included_df.copy()
    if "unit_source" not in base_df.columns:
        base_df["unit_source"] = "unit_quality_table"
    return pd.concat([base_df, pd.DataFrame(fallback_rows)], ignore_index=True)


def build_pptx(config: dict, logger: PipelineLogger) -> None:
    paths = resolve_project_paths(config)
    stim_df = normalize_stim_schedule(read_table(paths["stim_schedule_path"]), file_id_column=config["project"]["file_id_column"])
    unit_df = normalize_include_column(read_table(paths["unit_quality_path"]))
    file_id_column = config["project"]["file_id_column"]
    if "pl2_file" not in unit_df.columns:
        unit_df["pl2_file"] = ""
    stim_df[file_id_column] = [
        canonicalize_file_id(str(row[file_id_column]), row.get("pl2_file"), config)
        for row in stim_df.to_dict("records")
    ]
    unit_df[file_id_column] = [
        canonicalize_file_id(str(row[file_id_column]), row.get("pl2_file"), config)
        for row in unit_df.to_dict("records")
    ]
    included_df = unit_df[unit_df["include_bool"]]
    included_df = _augment_units_from_fullrate_fallback(config, paths, stim_df, included_df, logger)
    aligned_cfg = _aligned_cfg(config)
    generate_summary_figures(config=config, logger=logger)

    prs = Presentation()
    prs.slide_width = Inches(float(config["pptx"]["slide_width_in"]))
    prs.slide_height = Inches(float(config["pptx"]["slide_height_in"]))
    blank = prs.slide_layouts[6]

    layout = config["pptx"]["layout"]
    if layout == "one_unit_per_slide":
        for row in included_df.itertuples(index=False):
            raw_file_id = str(getattr(row, file_id_column))
            row_pl2_file = str(getattr(row, "pl2_file", "") or "")
            file_id = canonicalize_file_id(raw_file_id, row_pl2_file, config)
            unit_id = str(row.unit_id)
            original_name = str(row.original_name)
            note = str(row.note) if row.note is not None else ""
            unit_source = str(getattr(row, "unit_source", "unit_quality_table") or "unit_quality_table")
            stim_sub = stim_df[stim_df[file_id_column].astype(str) == file_id]
            if not stim_sub.empty:
                row_pl2_file = str(stim_sub["pl2_file"].iloc[0])
            inferred_has_light = infer_has_light_from_identifiers(raw_file_id, row_pl2_file)
            if stim_sub.empty:
                logger.log(
                    "build_pptx",
                    file_id,
                    row_pl2_file,
                    "",
                    "warning",
                    "No stim_schedule row matched this unit row after canonical file_id lookup; inferring has_light from identifiers.",
                )
            has_light = bool(stim_sub["has_light_bool"].any()) if not stim_sub.empty else (inferred_has_light if inferred_has_light is not None else True)
            event_group = str(stim_sub["event_group"].iloc[0]) if not stim_sub.empty and "event_group" in stim_sub.columns else ""
            duration_summary, light_on_summary, light_off_summary, n_events = _interval_metadata(config, paths, file_id, stim_sub)
            window_summary = _window_summary(config, stim_sub)

            full_png = _figure_path_with_legacy(
                paths=paths,
                config=config,
                file_id=file_id,
                pl2_file=row_pl2_file,
                unit_id=unit_id,
                figure_dir_key="figure_fullrate_dir",
                suffix="FullRate.png",
                logger=logger,
            )
            psth_png = paths["figure_psth_dir"] / f"{file_id}_{unit_id}_PSTH.png"
            aligned_png = paths["figure_aligned_dir"] / f"{file_id}_{unit_id}_AlignedRate_{aligned_window_tag(aligned_cfg)}.png"
            prepost_png = paths["figure_prepost_dir"] / f"{file_id}_{unit_id}_PreLightPost.png"
            if not has_light:
                aligned_png = _figure_path_with_legacy(
                    paths=paths,
                    config=config,
                    file_id=file_id,
                    pl2_file=row_pl2_file,
                    unit_id=unit_id,
                    figure_dir_key="figure_aligned_dir",
                    suffix="AlignedRate_no_light_skipped.png",
                    logger=logger,
                )
                prepost_png = _figure_path_with_legacy(
                    paths=paths,
                    config=config,
                    file_id=file_id,
                    pl2_file=row_pl2_file,
                    unit_id=unit_id,
                    figure_dir_key="figure_prepost_dir",
                    suffix="PreLightPost_no_light_skipped.png",
                    logger=logger,
                )
            analysis_mode = resolve_effective_analysis_mode(
                config,
                has_light=has_light,
                has_fullrate=full_png.exists(),
                has_aligned_assets=aligned_png.exists() or prepost_png.exists(),
            )

            slide = prs.slides.add_slide(blank)
            _add_textbox(slide, 0.3, 0.1, 12.7, 0.5, f"{file_id} | {unit_id} | {original_name}", font_size=22)

            if analysis_mode == "fullrate_aligned":
                _add_textbox(slide, 0.3, 0.65, 4.0, 0.25, "Panel A: Full-session rate", font_size=14)
                _add_textbox(slide, 4.55, 0.65, 4.0, 0.25, "Panel B: Light-aligned rate", font_size=14)
                _add_textbox(slide, 8.8, 0.65, 4.0, 0.25, "Panel C: Pre / light / post", font_size=14)
                if full_png.exists():
                    _add_picture_contain(slide, full_png, 0.3, 0.9, 4.0, 3.1)
                else:
                    _add_textbox(slide, 0.4, 1.8, 3.8, 1.0, f"Missing FullRate figure\n{full_png.name}", font_size=16)
                if aligned_png.exists():
                    _add_picture_contain(slide, aligned_png, 4.55, 0.9, 4.0, 3.1)
                elif not has_light:
                    _add_textbox(slide, 4.65, 1.8, 3.8, 1.0, "No light event.\nAligned analysis skipped.", font_size=16)
                else:
                    _add_textbox(slide, 4.65, 1.8, 3.8, 1.0, f"Missing aligned-rate figure\n{aligned_png.name}", font_size=16)
                if prepost_png.exists():
                    _add_picture_contain(slide, prepost_png, 8.8, 0.9, 4.0, 3.1)
                elif not has_light:
                    _add_textbox(slide, 8.9, 1.8, 3.8, 1.0, "No light event.\nPre/light/post summary not applicable.", font_size=16)
                else:
                    _add_textbox(slide, 8.9, 1.8, 3.8, 1.0, f"Missing pre/post summary\n{prepost_png.name}", font_size=16)
            else:
                if psth_png.exists():
                    _add_picture_contain(slide, psth_png, 0.3, 0.8, 6.1, 3.2)
                else:
                    _add_textbox(slide, 0.3, 1.5, 6.1, 1.0, f"Missing PSTH figure\n{psth_png.name}", font_size=18)
                    logger.log("build_pptx", file_id, "", str(psth_png), "warning", "PSTH figure missing; inserted placeholder text.")

                if full_png.exists():
                    _add_picture_contain(slide, full_png, 6.8, 0.8, 6.1, 3.2)
                else:
                    _add_textbox(slide, 6.8, 1.5, 6.1, 1.0, f"Missing FullRate figure\n{full_png.name}", font_size=18)
                    logger.log("build_pptx", file_id, "", str(full_png), "warning", "FullRate figure missing; inserted placeholder text.")

            if config["pptx"].get("include_metadata", True):
                meta_lines = _metadata_lines(
                    config,
                    file_id,
                    unit_id,
                    original_name,
                    note,
                    n_events,
                    duration_summary,
                    light_on_summary,
                    light_off_summary,
                    window_summary,
                    has_light=has_light,
                    event_group=event_group,
                    unit_source=unit_source,
                    effective_mode=analysis_mode,
                )
                _add_metadata_two_columns(slide, 0.35, 4.15, 12.2, 2.75, meta_lines, font_size=10)

    elif layout == "one_file_per_slide":
        for file_id, file_units in included_df.groupby(config["project"]["file_id_column"], sort=False):
            slide = prs.slides.add_slide(blank)
            _add_textbox(slide, 0.3, 0.1, 12.5, 0.5, f"{file_id} | File summary", font_size=22)

            summary_tagged = paths["figure_summary_dir"] / f"{file_id}_Summary_{aligned_window_tag(aligned_cfg)}.png"
            summary_plain = paths["figure_summary_dir"] / f"{file_id}_Summary.png"
            stim_sub = stim_df[stim_df[config["project"]["file_id_column"]].astype(str) == str(file_id)]
            has_light = bool(stim_sub["has_light_bool"].any()) if not stim_sub.empty else True
            summary_no_light = paths["figure_summary_dir"] / f"{file_id}_Summary_no_light.png"
            effective_mode = resolve_effective_analysis_mode(config, has_light=has_light, has_aligned_assets=summary_tagged.exists())
            summary_name = summary_no_light.name if not has_light else (summary_tagged.name if effective_mode == "fullrate_aligned" else summary_plain.name)
            summary_png = paths["figure_summary_dir"] / summary_name
            overlay_png = paths["figure_summary_dir"] / f"{file_id}_AllUnits_FullRate.png"
            included_units = ", ".join(file_units["unit_id"].astype(str).tolist())
            excluded_units = unit_df[
                (unit_df[config["project"]["file_id_column"]] == file_id) & (~unit_df["include_bool"])
            ]["unit_id"].astype(str).tolist()
            if has_light:
                stim_windows = ", ".join(f"{on}-{off}s" for on, off in zip(stim_sub["light_on_s"], stim_sub["light_off_s"]))
            else:
                stim_windows = "no light event"
            info = [
                f"included_units: {included_units or '-'}",
                f"excluded_or_duplicate_units: {', '.join(excluded_units) or '-'}",
                f"stim windows: {stim_windows}",
            ]
            _add_textbox(slide, 0.3, 5.7, 12.5, 1.4, "\n".join(info), font_size=14)

            if summary_png.exists():
                _add_picture_contain(slide, summary_png, 0.3, 0.8, 8.4, 4.6)
            else:
                _add_textbox(slide, 0.4, 1.8, 7.8, 1.0, f"Missing summary figure\n{summary_png.name}", font_size=18)
                logger.log("build_pptx", str(file_id), "", str(summary_png), "warning", "Summary figure missing; inserted placeholder text.")

            if overlay_png.exists():
                _add_picture_contain(slide, overlay_png, 8.9, 0.8, 3.9, 4.6)
            else:
                _add_textbox(slide, 9.0, 1.8, 3.6, 1.0, f"Missing overlay\n{overlay_png.name}", font_size=16)
                logger.log("build_pptx", str(file_id), "", str(overlay_png), "warning", "Overlay figure missing; inserted placeholder text.")

    else:
        for summary_png in sorted(paths["figure_summary_dir"].glob("*.png")):
            slide = prs.slides.add_slide(blank)
            _add_textbox(slide, 0.3, 0.1, 12.5, 0.5, summary_png.stem, font_size=22)
            _add_picture_contain(slide, summary_png, 0.4, 0.8, 12.2, 6.3)

    output_path = paths["pptx_output_path"]
    if config["run"]["dry_run"]:
        logger.log("build_pptx", "*", "", str(output_path), "skipped", "Dry-run mode: PPTX not written.")
        return

    try:
        prs.save(str(output_path))
        logger.log("build_pptx", "*", "", str(output_path), "success", "PPTX exported successfully.")
    except PermissionError as exc:
        fallback_output = output_path.with_name(
            f"{output_path.stem}_autosave_{datetime.now().strftime('%Y%m%d_%H%M%S')}{output_path.suffix}"
        )
        prs.save(str(fallback_output))
        logger.log(
            "build_pptx",
            "*",
            "",
            str(fallback_output),
            "warning",
            "Primary PPTX output was locked by another process; saved to fallback filename instead.",
            exception=exc,
        )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build PPTX summary slides from exported figures.")
    parser.add_argument("--config", required=True)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    config = load_yaml(Path(args.config))
    logger = PipelineLogger(resolve_project_paths(config)["logs_dir"])
    try:
        build_pptx(config=config, logger=logger)
        return 0
    finally:
        logger.save()


if __name__ == "__main__":
    raise SystemExit(main())
