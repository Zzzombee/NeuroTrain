from __future__ import annotations

import sys
import tempfile
import unittest
from pathlib import Path
from unittest import mock

import pandas as pd
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.build_aligned_rate_from_fullrate import build_aligned_rate_from_fullrate
from scripts.build_pptx import build_pptx
from scripts.build_stim_schedule_from_filenames import build_stim_schedule_from_filenames, parse_pl2_filename
from scripts.build_unit_quality_table import build_unit_quality_table
from scripts.export_figures import generate_summary_figures
from scripts.maintenance import canonicalize_project_tables
from utils.analysis_mode_utils import resolve_effective_analysis_mode
from utils.file_id_utils import canonicalize_file_id, canonicalize_file_id_from_pl2_file
from utils.logging_utils import PipelineLogger
from utils.path_utils import resolve_project_paths
from utils.table_utils import read_table, write_table
from tests.test_neuroexplorer_nex_backend import sample_config


class NoLightAutoTests(unittest.TestCase):
    def _config(self, tmp_root: Path) -> tuple[dict, PipelineLogger]:
        for relative_dir in [
            "00_raw_pl2",
            "01_sorting_info",
            "02_stim_events",
            "03_nex_exports/fullrate",
            "03_nex_exports/aligned_rate",
            "05_exported_figures/fullrate",
            "05_exported_figures/aligned_rate",
            "05_exported_figures/prepost_summary",
            "05_exported_figures/summary",
            "06_pptx",
            "99_logs",
        ]:
            (tmp_root / relative_dir).mkdir(parents=True, exist_ok=True)
        config = sample_config(tmp_root)
        config["analysis"]["mode"] = "auto"
        config["neuroexplorer"]["export_fullrate"] = True
        config["neuroexplorer"]["fullrate"]["enabled"] = True
        config["neuroexplorer"]["export_psth"] = False
        config["input"]["stim_schedule"] = "02_stim_events/stim_schedule_master.csv"
        config["stim_schedule"]["output_path"] = "02_stim_events/stim_schedule_master.csv"
        config["input"]["unit_quality_table"] = "01_sorting_info/unit_quality_table.csv"
        config["unit_table"]["output_path"] = "01_sorting_info/unit_quality_table.csv"
        config["unit_table"]["source"]["fallback_to_existing_fullrate_exports"] = True
        config["run"]["modules"] = {
            "validate": True,
            "build_stim_schedule": True,
            "prepare_events": True,
            "build_unit_table": True,
            "neuroexplorer_export": True,
            "aligned_rate": True,
            "origin_plot": True,
            "build_pptx": True,
        }
        logger = PipelineLogger(resolve_project_paths(config)["logs_dir"])
        return config, logger

    def _write_light_and_no_light_inputs(self, tmp_root: Path, config: dict) -> dict:
        paths = resolve_project_paths(config)
        (paths["pl2_dir"] / "sorted_01_120light15_1.pl2").write_text("", encoding="utf-8")
        (paths["pl2_dir"] / "sorted_02_nolight_1.pl2").write_text("", encoding="utf-8")
        write_table(
            pd.DataFrame(
                [
                    {
                        "file_id": "01",
                        "pl2_file": "sorted_01_120light15_1.pl2",
                        "event_group": "120light15",
                        "has_light": "yes",
                        "light_on_s": 120,
                        "duration_s": 15,
                        "light_off_s": 135,
                    },
                    {
                        "file_id": "02",
                        "pl2_file": "sorted_02_nolight_1.pl2",
                        "event_group": "nolight",
                        "has_light": "no",
                        "light_on_s": "",
                        "duration_s": "",
                        "light_off_s": "",
                        "condition": "no_light",
                    },
                ]
            ),
            paths["stim_schedule_path"],
        )
        write_table(
            pd.DataFrame(
                [
                    {"file_id": "01", "unit_id": "unit01", "original_name": "unit01", "include": "yes"},
                    {"file_id": "02", "unit_id": "unit01", "original_name": "unit01", "include": "yes"},
                ]
            ),
            paths["unit_quality_path"],
        )
        fullrate_df = pd.DataFrame(
            {
                "file_id": ["01"] * 4 + ["02"] * 4,
                "unit_id": ["unit01"] * 8,
                "time_bin_start_s": [59.5, 119.5, 134.5, 194.5, 0.5, 1.5, 2.5, 3.5],
                "time_bin_end_s": [60.5, 120.5, 135.5, 195.5, 1.5, 2.5, 3.5, 4.5],
                "time_bin_center_s": [60, 120, 135, 195, 1, 2, 3, 4],
                "firing_rate_hz": [1, 3, 2, 1, 4, 5, 4, 6],
            }
        )
        write_table(fullrate_df[fullrate_df["file_id"] == "01"], paths["nex_fullrate_dir"] / "01_FullRate_bin1s.csv")
        write_table(fullrate_df[fullrate_df["file_id"] == "02"], paths["nex_fullrate_dir"] / "02_FullRate_bin1s.csv")
        return paths

    def test_resolve_auto_no_light_is_fullrate_aligned(self):
        config = sample_config(Path("D:/tmp"))
        config["analysis"]["mode"] = "auto"
        self.assertEqual(resolve_effective_analysis_mode(config, has_light=False), "fullrate_aligned")

    def test_canonicalize_file_id(self):
        config = sample_config(Path("D:/tmp"))
        self.assertEqual(canonicalize_file_id_from_pl2_file("sorted_01_nolight_1.pl2", config), "01")
        self.assertEqual(canonicalize_file_id("test01", "sorted_01_nolight_1.pl2", config), "01")
        self.assertEqual(canonicalize_file_id("sorted_01_nolight_1", None, config), "01")

    def test_build_stim_schedule_parses_no_light(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._config(tmp_root)
            parsed = parse_pl2_filename(config, Path("sorted_02_nolight_1,5,9.pl2"))
            self.assertEqual(parsed["file_id"], "02")
            self.assertEqual(parsed["has_light"], "no")
            self.assertEqual(parsed["event_group"], "nolight")
            self.assertEqual(parsed["light_on_s"], "")

            (resolve_project_paths(config)["pl2_dir"] / "sorted_02_nolight_1,5,9.pl2").write_text("", encoding="utf-8")
            output = build_stim_schedule_from_filenames(config, logger)
            df = read_table(output)
            self.assertEqual(df.loc[0, "has_light"], "no")
            self.assertTrue(pd.isna(df.loc[0, "light_on_s"]))

    def test_no_light_aligned_status_and_light_file_outputs(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._config(tmp_root)
            paths = self._write_light_and_no_light_inputs(tmp_root, config)
            build_aligned_rate_from_fullrate(config, logger)
            self.assertTrue((paths["nex_aligned_rate_dir"] / "01_LightAlignedRate_pre60_post85_bin1s.csv").exists())
            self.assertTrue((paths["nex_aligned_rate_dir"] / "02_LightAlignedRate_no_light_skipped.csv").exists())
            status_df = read_table(paths["nex_aligned_rate_dir"] / "02_PreLightPostSummary_no_light_skipped.csv")
            self.assertEqual(status_df.loc[0, "analysis_status"], "no_light_skipped")

    def test_export_figures_and_pptx_no_light_page(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._config(tmp_root)
            paths = self._write_light_and_no_light_inputs(tmp_root, config)
            build_aligned_rate_from_fullrate(config, logger)
            full_png = paths["figure_fullrate_dir"] / "02_unit01_FullRate.png"
            fig_df = pd.DataFrame({"x": [0, 1], "y": [1, 2]})
            import matplotlib.pyplot as plt

            fig, ax = plt.subplots()
            ax.plot(fig_df["x"], fig_df["y"])
            fig.savefig(full_png)
            fig.savefig(paths["figure_fullrate_dir"] / "01_unit01_FullRate.png")
            plt.close(fig)

            generate_summary_figures(config, logger)
            self.assertTrue((paths["figure_aligned_dir"] / "02_unit01_AlignedRate_no_light_skipped.png").exists())
            self.assertTrue((paths["figure_prepost_dir"] / "02_unit01_PreLightPost_no_light_skipped.png").exists())
            self.assertTrue((paths["figure_summary_dir"] / "02_Summary_no_light.png").exists())

            build_pptx(config, logger)
            pptx_path = paths["pptx_output_path"]
            prs = Presentation(str(pptx_path))
            all_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
            self.assertIn("has_light: no", all_text)
            self.assertIn("event_group: nolight", all_text)
            self.assertIn("alignment: not applicable", all_text)
            self.assertNotIn("Missing FullRate figure\x0b02_unit01_FullRate.png", all_text)
            metadata_shapes = [
                shape
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and ("file_id:" in shape.text or "window_mode:" in shape.text or "alignment:" in shape.text)
            ]
            self.assertGreaterEqual(len(metadata_shapes), 2)
            for shape in metadata_shapes:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            self.assertLessEqual(run.font.size.pt, 10)

    def test_unit_table_falls_back_to_existing_fullrate_csv(self):
        class FailingAdapter:
            def __init__(self, config, logger): pass
            def connect(self): raise RuntimeError("nex unavailable")
            def close_file(self): pass
            def quit(self): pass

        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._config(tmp_root)
            paths = resolve_project_paths(config)
            (paths["pl2_dir"] / "sorted_02_nolight_1.pl2").write_text("", encoding="utf-8")
            write_table(
                pd.DataFrame(
                    [
                        {
                            "file_id": "02",
                            "pl2_file": "sorted_02_nolight_1.pl2",
                            "event_group": "nolight",
                            "has_light": "no",
                            "light_on_s": "",
                            "duration_s": "",
                        }
                    ]
                ),
                paths["stim_schedule_path"],
            )
            write_table(
                pd.DataFrame(
                    {
                        "file_id": ["02", "02"],
                        "unit_id": ["unitA", "unitB"],
                        "time_bin_center_s": [1, 1],
                        "firing_rate_hz": [2, 3],
                    }
                ),
                paths["nex_fullrate_dir"] / "02_FullRate_bin1s.csv",
            )
            with mock.patch("scripts.build_unit_quality_table.NeuroExplorerAdapter", side_effect=lambda config, logger: FailingAdapter(config, logger)):
                build_unit_quality_table(config, logger)
            df = read_table(paths["unit_quality_path"])
            self.assertEqual(df["detected_by"].unique().tolist(), ["existing_fullrate_export_fallback"])
            self.assertEqual(df["include"].tolist(), ["yes", "yes"])

    def test_canonicalize_merges_duplicate_unit_rows_and_preserves_manual_fields(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._config(tmp_root)
            paths = resolve_project_paths(config)
            write_table(
                pd.DataFrame(
                    [
                        {
                            "file_id": "test01",
                            "pl2_file": "sorted_01_nolight_1.pl2",
                            "unit_id": "unit01",
                            "original_name": "SPK_SPKC01a",
                            "include": "no",
                            "exclusion_reason": "manual exclusion",
                            "duplicate_of": "unit02",
                            "representative_unit": "unit02",
                            "note": "keep me",
                        },
                        {
                            "file_id": "sorted_01_nolight_1",
                            "pl2_file": "sorted_01_nolight_1.pl2",
                            "unit_id": "unit01",
                            "original_name": "SPK_SPKC01a",
                            "include": "yes",
                            "exclusion_reason": "",
                            "duplicate_of": "",
                            "representative_unit": "",
                            "note": "",
                        },
                    ]
                ),
                paths["unit_quality_path"],
            )
            canonicalize_project_tables(config, logger)
            df = read_table(paths["unit_quality_path"])
            self.assertEqual(len(df), 1)
            self.assertEqual(str(df.loc[0, "file_id"]), "01")
            self.assertEqual(str(df.loc[0, "include"]), "no")
            self.assertEqual(str(df.loc[0, "exclusion_reason"]), "manual exclusion")
            self.assertEqual(str(df.loc[0, "duplicate_of"]), "unit02")
            self.assertEqual(str(df.loc[0, "note"]), "keep me")

    def test_build_pptx_uses_legacy_fullrate_only_as_fallback(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._config(tmp_root)
            paths = self._write_light_and_no_light_inputs(tmp_root, config)
            legacy_png = paths["figure_fullrate_dir"] / "test02_unit01_FullRate.png"
            import matplotlib.pyplot as plt

            fig, ax = plt.subplots()
            ax.plot([0, 1], [1, 2])
            fig.savefig(legacy_png)
            plt.close(fig)
            generate_summary_figures(config, logger)
            build_pptx(config, logger)
            self.assertTrue(any("Using legacy file_id figure path" in record.message for record in logger.records))


if __name__ == "__main__":
    unittest.main()
