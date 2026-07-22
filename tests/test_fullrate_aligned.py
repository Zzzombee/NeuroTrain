from __future__ import annotations

import sys
import tempfile
import unittest
from pathlib import Path
from unittest import mock

import pandas as pd
import numpy as np

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.build_aligned_rate_from_fullrate import build_aligned_rate_for_file, build_aligned_rate_from_fullrate
from scripts.build_time_cluster_aligned_rate import build_time_cluster_aligned_rate_for_file
from build_pptx import build_pptx
from scripts.export_from_neuroexplorer import export_from_neuroexplorer
from plot_in_origin import plot_in_origin
from utils.aligned_utils import compute_aligned_window, compute_pre_light_post_windows
from utils.logging_utils import PipelineLogger
from utils.path_utils import resolve_project_paths
from utils.table_utils import write_table, read_table
from validate_project import validate_project
from tests.test_neuroexplorer_nex_backend import sample_config


class FullrateAlignedTests(unittest.TestCase):
    def _init_project(self, tmp_root: Path) -> tuple[dict, PipelineLogger]:
        for relative_dir in [
            "00_raw_pl2",
            "01_sorting_info",
            "02_stim_events",
            "03_nex_exports/fullrate",
            "03_nex_exports/aligned_rate",
            "99_logs",
        ]:
            (tmp_root / relative_dir).mkdir(parents=True, exist_ok=True)
        (tmp_root / "00_raw_pl2" / "demo.pl2").write_text("", encoding="utf-8")
        (tmp_root / "02_stim_events" / "stim_schedule_master.csv").write_text(
            "file_id,pl2_file,has_light,light_on_s,duration_s,light_off_s\n"
            "demo,demo.pl2,yes,120,15,135\n"
            "demo,demo.pl2,yes,240,15,255\n",
            encoding="utf-8",
        )
        (tmp_root / "01_sorting_info" / "unit_quality_table.csv").write_text(
            "file_id,unit_id,original_name,include\n"
            "demo,unit01,SPK_SPKC04a,yes\n",
            encoding="utf-8",
        )
        config = sample_config(tmp_root)
        config["analysis"]["mode"] = "fullrate_aligned"
        config["neuroexplorer"]["fullrate"]["enabled"] = True
        config["neuroexplorer"]["export_fullrate"] = True
        config["aligned_rate"]["window_mode"] = "light_duration_plus_margin"
        config["aligned_rate"]["pre_margin_s"] = 60
        config["aligned_rate"]["post_margin_s"] = 60
        logger = PipelineLogger(resolve_project_paths(config)["logs_dir"])
        return config, logger

    def test_build_aligned_rate_for_file(self):
        config = sample_config(Path("D:/tmp"))
        config["analysis"]["mode"] = "fullrate_aligned"
        fullrate_df = pd.DataFrame(
            {
                "file_id": ["demo"] * 8,
                "unit_id": ["unit01"] * 8,
                "time_bin_start_s": [59.5, 104.5, 119.5, 134.5, 194.5, 239.5, 254.5, 314.5],
                "time_bin_end_s": [60.5, 105.5, 120.5, 135.5, 195.5, 240.5, 255.5, 315.5],
                "time_bin_center_s": [60, 105, 120, 135, 195, 240, 255, 315],
                "firing_rate_hz": [1, 2, 4, 3, 2, 5, 4, 2],
                "source_file": ["raw.txt"] * 8,
            }
        )
        stim_sub = pd.DataFrame(
            {
                "file_id": ["demo", "demo"],
                "pl2_file": ["demo.pl2", "demo.pl2"],
                "light_on_s": [120.0, 240.0],
                "duration_s": [15.0, 15.0],
                "light_off_s": [135.0, 255.0],
            }
        )
        aligned_df, summary_df = build_aligned_rate_for_file(config, "demo", fullrate_df, stim_sub)
        self.assertIn("aligned_time_s", aligned_df.columns)
        self.assertEqual(
            sorted(aligned_df["aligned_time_s"].unique().tolist()),
            [-60.0, -45.0, -15.0, 0.0, 15.0, 75.0],
        )
        self.assertIn(0.0, aligned_df["aligned_time_s"].tolist())
        self.assertNotIn("aligned_bin_start_s", aligned_df.columns)
        self.assertEqual(aligned_df["aligned_x_min_s"].iloc[0], -60.0)
        self.assertEqual(aligned_df["aligned_x_max_s"].iloc[0], 85.0)
        aggregated = summary_df[summary_df["trial_id"].astype(str) == "aggregated"]
        self.assertFalse(aggregated.empty)
        trial_rows = summary_df[summary_df["trial_id"].astype(str) == "1"]
        self.assertEqual(trial_rows["baseline_window_start_s"].iloc[0], -60.0)
        self.assertEqual(trial_rows["baseline_window_end_s"].iloc[0], 0.0)
        self.assertEqual(trial_rows["light_window_start_s"].iloc[0], 5.0)
        self.assertEqual(trial_rows["light_window_end_s"].iloc[0], 20.0)
        self.assertEqual(trial_rows["post_window_start_s"].iloc[0], 25.0)
        self.assertEqual(trial_rows["post_window_end_s"].iloc[0], 85.0)
        self.assertEqual(trial_rows["summary_window_mode"].iloc[0], "configured_windows")

    def test_time_cluster_ten_second_bins_use_zero_as_boundary(self):
        config = sample_config(Path("D:/tmp"))
        config["time_cluster_aligned_rate"] = {
            "window_s": [-120, 320],
            "source_bin_width_s": 10,
            "bin_width_s": 10,
            "require_light_on_on_bin_boundary": True,
            "off_boundary_policy": "error",
        }
        starts = np.arange(0.0, 450.0, 10.0)
        fullrate_df = pd.DataFrame(
            {
                "file_id": "demo",
                "unit_id": "unit01",
                "time_bin_start_s": starts,
                "time_bin_end_s": starts + 10.0,
                "time_bin_center_s": starts + 5.0,
                "firing_rate_hz": np.arange(len(starts), dtype=float),
            }
        )
        stim_sub = pd.DataFrame(
            {
                "file_id": ["demo"],
                "light_on_s": [120.0],
                "duration_s": [25.0],
                "light_off_s": [145.0],
            }
        )
        aligned_df = build_time_cluster_aligned_rate_for_file(config, "demo", fullrate_df, stim_sub)
        centers = aligned_df["aligned_time_s"].to_numpy(dtype=float)
        np.testing.assert_array_equal(centers[:3], np.asarray([-115.0, -105.0, -95.0]))
        self.assertIn(-5.0, centers)
        self.assertIn(5.0, centers)
        self.assertIn(15.0, centers)
        self.assertNotIn(0.0, centers)
        baseline = centers[(centers >= -120.0) & (centers < 0.0)]
        test = centers[(centers >= 0.0) & (centers < 300.0)]
        np.testing.assert_array_equal(baseline, np.arange(-115.0, 0.0, 10.0))
        np.testing.assert_array_equal(test, np.arange(5.0, 300.0, 10.0))
        self.assertEqual(len(baseline), 12)
        self.assertEqual(len(test), 30)
        self.assertTrue((aligned_df["aligned_bin_start_s"] % 10.0 == 0.0).all())
        self.assertTrue(aligned_df["alignment_exact"].all())

    def test_compute_pre_light_post_windows_configured(self):
        config = sample_config(Path("D:/tmp"))
        for duration, expected in [
            (15.0, (-60.0, 0.0, 5.0, 20.0, 25.0, 85.0)),
            (25.0, (-60.0, 0.0, 5.0, 20.0, 25.0, 85.0)),
            (60.0, (-60.0, 0.0, 5.0, 20.0, 25.0, 85.0)),
        ]:
            windows = compute_pre_light_post_windows(duration, config["aligned_rate"])
            self.assertEqual(windows["summary_window_mode"], "configured_windows")
            self.assertEqual(
                (
                    windows["baseline_window_start_s"],
                    windows["baseline_window_end_s"],
                    windows["light_window_start_s"],
                    windows["light_window_end_s"],
                    windows["post_window_start_s"],
                    windows["post_window_end_s"],
                ),
                expected,
            )

    def test_compute_pre_light_post_windows_fixed_mode(self):
        config = sample_config(Path("D:/tmp"))
        config["aligned_rate"].pop("pre_window_s", None)
        config["aligned_rate"].pop("light_window_s", None)
        config["aligned_rate"].pop("post_window_s", None)
        config["aligned_rate"]["summary_window_mode"] = "fixed"
        config["aligned_rate"]["baseline_window_s"] = [-15, 0]
        config["aligned_rate"]["post_window_s"] = [15, 30]
        windows = compute_pre_light_post_windows(25.0, config["aligned_rate"])
        self.assertEqual(windows["summary_window_mode"], "fixed")
        self.assertEqual(windows["baseline_window_start_s"], -15.0)
        self.assertEqual(windows["post_window_end_s"], 30.0)

    def test_compute_aligned_window_dynamic_mode(self):
        config = sample_config(Path("D:/tmp"))
        window = compute_aligned_window(200.0, 225.0, 25.0, config["aligned_rate"])
        self.assertEqual(window["abs_start_s"], 140.0)
        self.assertEqual(window["abs_end_s"], 285.0)
        self.assertEqual(window["aligned_x_min_s"], -60.0)
        self.assertEqual(window["aligned_x_max_s"], 85.0)

    def test_compute_aligned_window_uses_duration_when_light_off_missing(self):
        config = sample_config(Path("D:/tmp"))
        window = compute_aligned_window(200.0, None, 25.0, config["aligned_rate"])
        self.assertEqual(window["light_off_s"], 225.0)
        self.assertEqual(window["aligned_x_max_s"], 85.0)

    def test_compute_aligned_window_duration_60(self):
        config = sample_config(Path("D:/tmp"))
        window = compute_aligned_window(200.0, 260.0, 60.0, config["aligned_rate"])
        self.assertEqual(window["aligned_x_max_s"], 85.0)
        self.assertEqual(window["abs_end_s"], 285.0)

    def test_compute_aligned_window_fixed_mode_compatible(self):
        config = sample_config(Path("D:/tmp"))
        config["aligned_rate"].pop("pre_window_s", None)
        config["aligned_rate"].pop("light_window_s", None)
        config["aligned_rate"].pop("post_window_s", None)
        config["aligned_rate"]["window_mode"] = "fixed"
        config["aligned_rate"]["x_min_s"] = -60
        config["aligned_rate"]["x_max_s"] = 75
        window = compute_aligned_window(120.0, 135.0, 15.0, config["aligned_rate"])
        self.assertEqual(window["abs_start_s"], 60.0)
        self.assertEqual(window["abs_end_s"], 195.0)
        self.assertEqual(window["aligned_x_min_s"], -60.0)
        self.assertEqual(window["aligned_x_max_s"], 75.0)

    def test_build_aligned_rate_from_fullrate_writes_outputs(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            paths = resolve_project_paths(config)
            fullrate_df = pd.DataFrame(
                {
                    "file_id": ["demo"] * 8,
                    "unit_id": ["SPK_SPKC04a"] * 8,
                    "time_bin_start_s": [59.5, 104.5, 119.5, 134.5, 194.5, 239.5, 254.5, 314.5],
                    "time_bin_end_s": [60.5, 105.5, 120.5, 135.5, 195.5, 240.5, 255.5, 315.5],
                    "time_bin_center_s": [60, 105, 120, 135, 195, 240, 255, 315],
                    "firing_rate_hz": [1, 2, 4, 3, 2, 5, 4, 2],
                    "source_file": ["raw.txt"] * 8,
                }
            )
            write_table(fullrate_df, paths["nex_fullrate_dir"] / "demo_FullRate_bin1s.csv")
            build_aligned_rate_from_fullrate(config, logger)
            aligned_path = paths["nex_aligned_rate_dir"] / "demo_LightAlignedRate_pre60_post85_bin1s.csv"
            summary_path = paths["nex_aligned_rate_dir"] / "demo_PreLightPostSummary.csv"
            self.assertTrue(aligned_path.exists())
            self.assertTrue(summary_path.exists())
            self.assertIn("pre60_post85", aligned_path.name)

    def test_build_aligned_rate_warns_when_window_precedes_recording_start(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            (tmp_root / "02_stim_events" / "stim_schedule_master.csv").write_text(
                "file_id,pl2_file,has_light,light_on_s,duration_s,light_off_s\n"
                "demo,demo.pl2,yes,30,25,55\n",
                encoding="utf-8",
            )
            paths = resolve_project_paths(config)
            write_table(
                pd.DataFrame(
                    {
                        "file_id": ["demo"] * 3,
                        "unit_id": ["SPK_SPKC04a"] * 3,
                        "time_bin_start_s": [0.5, 29.5, 54.5],
                        "time_bin_end_s": [1.5, 30.5, 55.5],
                        "time_bin_center_s": [1, 30, 55],
                        "firing_rate_hz": [1, 2, 3],
                        "source_file": ["raw.txt"] * 3,
                    }
                ),
                paths["nex_fullrate_dir"] / "demo_FullRate_bin1s.csv",
            )
            build_aligned_rate_from_fullrate(config, logger)
            self.assertTrue(any("extends before recording start" in record.message for record in logger.records))

    def test_build_aligned_rate_skips_no_light_file(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            (tmp_root / "02_stim_events" / "stim_schedule_master.csv").write_text(
                "file_id,pl2_file,has_light,light_on_s,duration_s,light_off_s,event_group,condition\n"
                "demo,demo.pl2,no,,,,nolight,no_light\n",
                encoding="utf-8",
            )
            paths = resolve_project_paths(config)
            write_table(
                pd.DataFrame(
                    {
                        "file_id": ["demo"] * 2,
                        "unit_id": ["SPK_SPKC04a"] * 2,
                        "time_bin_start_s": [0.5, 1.5],
                        "time_bin_end_s": [1.5, 2.5],
                        "time_bin_center_s": [1, 2],
                        "firing_rate_hz": [1, 2],
                        "source_file": ["raw.txt"] * 2,
                    }
                ),
                paths["nex_fullrate_dir"] / "demo_FullRate_bin1s.csv",
            )
            build_aligned_rate_from_fullrate(config, logger)
            self.assertTrue(any(record.status == "skipped" and "No light event; aligned rate skipped." in record.message for record in logger.records))

    def test_validate_project_allows_no_light_rows(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            (tmp_root / "02_stim_events" / "stim_schedule_master.csv").write_text(
                "file_id,pl2_file,has_light,event_group,condition,light_on_s,duration_s,light_off_s\n"
                "demo,demo.pl2,no,nolight,no_light,,,\n",
                encoding="utf-8",
            )
            validate_project(config, logger)
            self.assertTrue(any("Validated no-light control row" in record.message for record in logger.records))

    def test_validate_project_uses_configured_windows_without_duration_margin_warning(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            (tmp_root / "02_stim_events" / "stim_schedule_master.csv").write_text(
                "file_id,pl2_file,has_light,light_on_s,duration_s,light_off_s\n"
                "demo,demo.pl2,yes,200,120,320\n",
                encoding="utf-8",
            )
            validate_project(config, logger)
            self.assertFalse(any("baseline window exceeds aligned pre-margin" in record.message for record in logger.records))

    def test_fullrate_aligned_export_does_not_require_events(self):
        class FakeAdapter:
            ensure_events_called = False

            def __init__(self, config, logger):
                self.config = config
                self.logger = logger
                self.paths = resolve_project_paths(config)

            def connect(self): pass
            def smoke_test(self): pass
            def introspect(self): pass
            def open_file(self, pl2_path): pass
            def ensure_events(self, *args, **kwargs): FakeAdapter.ensure_events_called = True
            def validate_required_events(self): raise AssertionError("validate_required_events should not run in fullrate_aligned")
            def get_reference_event_name(self): return "Light_On"
            def configure_psth_template(self, *args): raise AssertionError("PSTH should not run in fullrate_aligned")
            def run_template(self, template_name): self.last_template = template_name
            def export_psth(self, *args): raise AssertionError("PSTH export should not run in fullrate_aligned")
            def configure_fullrate_template(self, *args): pass
            def export_fullrate(self, file_id, unit_names, output_csv):
                output_csv.parent.mkdir(parents=True, exist_ok=True)
                output_csv.write_text("file_id,unit_id,time_bin_start_s,time_bin_end_s,time_bin_center_s,firing_rate_hz,source_file\n", encoding="utf-8")
            def close_file(self): pass
            def quit(self): pass

        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            with mock.patch("scripts.export_from_neuroexplorer.NeuroExplorerAdapter", side_effect=lambda config, logger: FakeAdapter(config, logger)):
                export_from_neuroexplorer(config, logger)
            messages = [record.message for record in logger.records]
            self.assertFalse(FakeAdapter.ensure_events_called)
            self.assertTrue(any("analysis_mode=fullrate_aligned" in message for message in messages))
            self.assertTrue(any("preferring fullrate_aligned export path" in message for message in messages))

    def test_auto_prefers_fullrate_aligned_before_psth(self):
        class FakeAdapter:
            ensure_events_called = False
            fullrate_run_called = False

            def __init__(self, config, logger):
                self.config = config
                self.logger = logger

            def connect(self): pass
            def smoke_test(self): pass
            def introspect(self): pass
            def open_file(self, pl2_path): pass
            def ensure_events(self, *args, **kwargs): FakeAdapter.ensure_events_called = True
            def validate_required_events(self): FakeAdapter.ensure_events_called = True
            def get_reference_event_name(self): return "Light_On"
            def configure_psth_template(self, *args): FakeAdapter.ensure_events_called = True
            def run_template(self, template_name):
                if template_name == "RateHist_FullSession":
                    FakeAdapter.fullrate_run_called = True
            def export_psth(self, *args): FakeAdapter.ensure_events_called = True
            def configure_fullrate_template(self, *args): pass
            def export_fullrate(self, file_id, unit_names, output_csv):
                output_csv.parent.mkdir(parents=True, exist_ok=True)
                output_csv.write_text("file_id,unit_id,time_bin_start_s,time_bin_end_s,time_bin_center_s,firing_rate_hz,source_file\n", encoding="utf-8")
            def close_file(self): pass
            def quit(self): pass

        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            config["analysis"]["mode"] = "auto"
            with mock.patch("scripts.export_from_neuroexplorer.NeuroExplorerAdapter", side_effect=lambda config, logger: FakeAdapter(config, logger)):
                export_from_neuroexplorer(config, logger)
            self.assertTrue(FakeAdapter.fullrate_run_called)
            self.assertFalse(FakeAdapter.ensure_events_called)

    def test_plot_in_origin_uses_expected_band_ranges(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            paths = resolve_project_paths(config)
            write_table(
                pd.DataFrame(
                    {
                        "file_id": ["demo"] * 4,
                        "unit_id": ["unit01"] * 4,
                        "time_bin_start_s": [59.5, 119.5, 134.5, 194.5],
                        "time_bin_end_s": [60.5, 120.5, 135.5, 195.5],
                        "time_bin_center_s": [60, 120, 135, 195],
                        "firing_rate_hz": [1, 2, 3, 2],
                        "source_file": ["raw.txt"] * 4,
                    }
                ),
                paths["nex_fullrate_dir"] / "demo_FullRate_bin1s.csv",
            )
            write_table(
                pd.DataFrame(
                    {
                        "file_id": ["demo"] * 4,
                        "unit_id": ["unit01"] * 4,
                        "trial_id": ["aggregated"] * 4,
                        "light_on_s": [120.0] * 4,
                        "light_off_s": [135.0] * 4,
                        "duration_s": [15.0] * 4,
                        "aligned_time_s": [-60.0, 0.0, 15.0, 75.0],
                        "firing_rate_hz": [1.0, 4.0, 3.0, 2.0],
                        "aggregation": ["mean"] * 4,
                    }
                ),
                paths["nex_aligned_rate_dir"] / "demo_LightAlignedRate_pre60_post85_bin1s.csv",
            )
            write_table(
                pd.DataFrame(
                    {
                        "file_id": ["demo"],
                        "unit_id": ["unit01"],
                        "trial_id": ["aggregated"],
                        "baseline_hz": [1.0],
                        "light_hz": [4.0],
                        "post_hz": [2.0],
                        "delta_light_minus_baseline": [3.0],
                        "ratio_light_to_baseline": [4.0],
                        "duration_s": [15.0],
                        "baseline_window_start_s": [-60.0],
                        "baseline_window_end_s": [0.0],
                        "light_window_start_s": [5.0],
                        "light_window_end_s": [20.0],
                        "post_window_start_s": [25.0],
                        "post_window_end_s": [85.0],
                        "summary_window_mode": ["configured_windows"],
                        "aggregation": ["mean"],
                    }
                ),
                paths["nex_aligned_rate_dir"] / "demo_PreLightPostSummary.csv",
            )

            captured = []

            def capture_export(**kwargs):
                captured.append((kwargs["title"], kwargs["band_ranges"]))

            with mock.patch("plot_in_origin._matplotlib_export", side_effect=capture_export), mock.patch(
                "plot_in_origin._summary_bar_export", return_value=None
            ), mock.patch("plot_in_origin.generate_summary_figures", return_value=None):
                plot_in_origin(config, logger)

            full_band = next(bands for title, bands in captured if "Full-session" in title)
            aligned_band = next(bands for title, bands in captured if "Light-aligned rate" in title)
            self.assertEqual(full_band, [(120.0, 135.0), (240.0, 255.0)])
            self.assertEqual(aligned_band, [(0.0, 15.0)])

    def test_plot_in_origin_overlay_matches_original_unit_names(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            paths = resolve_project_paths(config)
            write_table(
                pd.DataFrame(
                    {
                        "file_id": ["demo"] * 2,
                        "unit_id": ["SPK_SPKC04a"] * 2,
                        "time_bin_start_s": [0.5, 1.5],
                        "time_bin_end_s": [1.5, 2.5],
                        "time_bin_center_s": [1, 2],
                        "firing_rate_hz": [1, 2],
                        "source_file": ["raw.txt"] * 2,
                    }
                ),
                paths["nex_fullrate_dir"] / "demo_FullRate_bin1s.csv",
            )

            captured_overlay = []

            def capture_overlay(**kwargs):
                captured_overlay.append(kwargs["data"].copy())

            with mock.patch("plot_in_origin._matplotlib_export", return_value=None), mock.patch(
                "plot_in_origin._matplotlib_overlay_export", side_effect=capture_overlay
            ), mock.patch("plot_in_origin.generate_summary_figures", return_value=None):
                plot_in_origin(config, logger)

            self.assertEqual(len(captured_overlay), 1)
            self.assertEqual(captured_overlay[0]["unit_id"].unique().tolist(), ["SPK_SPKC04a"])

    def test_plot_in_origin_no_light_uses_no_full_band_and_placeholders(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            (tmp_root / "02_stim_events" / "stim_schedule_master.csv").write_text(
                "file_id,pl2_file,has_light,event_group,condition,light_on_s,duration_s,light_off_s\n"
                "demo,demo.pl2,no,nolight,no_light,,,\n",
                encoding="utf-8",
            )
            paths = resolve_project_paths(config)
            write_table(
                pd.DataFrame(
                    {
                        "file_id": ["demo"] * 2,
                        "unit_id": ["unit01"] * 2,
                        "time_bin_start_s": [0.5, 1.5],
                        "time_bin_end_s": [1.5, 2.5],
                        "time_bin_center_s": [1, 2],
                        "firing_rate_hz": [1, 2],
                        "source_file": ["raw.txt"] * 2,
                    }
                ),
                paths["nex_fullrate_dir"] / "demo_FullRate_bin1s.csv",
            )

            captured = []

            def capture_export(**kwargs):
                captured.append((kwargs["title"], kwargs["band_ranges"]))

            with mock.patch("plot_in_origin._matplotlib_export", side_effect=capture_export), mock.patch(
                "plot_in_origin.generate_summary_figures", return_value=None
            ):
                plot_in_origin(config, logger)

            full_band = next(bands for title, bands in captured if "Full-session" in title)
            self.assertEqual(full_band, [])
            self.assertTrue((paths["figure_aligned_dir"] / "demo_unit01_AlignedRate_no_light_skipped.png").exists())
            self.assertTrue((paths["figure_prepost_dir"] / "demo_unit01_PreLightPost_no_light_skipped.png").exists())

    def test_build_pptx_no_light_metadata(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            (tmp_root / "02_stim_events" / "stim_schedule_master.csv").write_text(
                "file_id,pl2_file,has_light,event_group,condition,note,light_on_s,duration_s,light_off_s\n"
                "demo,demo.pl2,no,nolight,no_light,\"sorted channels: 1,5,9\",,,\n",
                encoding="utf-8",
            )
            paths = resolve_project_paths(config)
            # create minimal figures expected by pptx builder
            import matplotlib.pyplot as plt
            for output_path in [
                paths["figure_fullrate_dir"] / "demo_unit01_FullRate.png",
                paths["figure_aligned_dir"] / "demo_unit01_AlignedRate_no_light_skipped.png",
                paths["figure_prepost_dir"] / "demo_unit01_PreLightPost_no_light_skipped.png",
                paths["figure_summary_dir"] / "demo_Summary_no_light.png",
            ]:
                fig, ax = plt.subplots()
                ax.text(0.5, 0.5, output_path.stem, ha="center", va="center")
                ax.axis("off")
                fig.savefig(output_path)
                plt.close(fig)
            with mock.patch("scripts.build_pptx.generate_summary_figures", return_value=None):
                build_pptx(config, logger)
            self.assertTrue(paths["pptx_output_path"].exists())
            from pptx import Presentation

            prs = Presentation(str(paths["pptx_output_path"]))
            slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
            self.assertIn("has_light: no", slide_text)
            self.assertIn("alignment: not applicable", slide_text)

    def test_build_pptx_metadata_uses_configured_summary_windows(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_root = Path(tmpdir)
            config, logger = self._init_project(tmp_root)
            (tmp_root / "02_stim_events" / "stim_schedule_master.csv").write_text(
                "file_id,pl2_file,has_light,light_on_s,duration_s,light_off_s\n"
                "demo,demo.pl2,yes,200,25,225\n",
                encoding="utf-8",
            )
            paths = resolve_project_paths(config)
            import matplotlib.pyplot as plt
            for output_path in [
                paths["figure_fullrate_dir"] / "demo_unit01_FullRate.png",
                paths["figure_aligned_dir"] / "demo_unit01_AlignedRate_pre60_post85.png",
                paths["figure_prepost_dir"] / "demo_unit01_PreLightPost.png",
            ]:
                fig, ax = plt.subplots()
                ax.text(0.5, 0.5, output_path.stem, ha="center", va="center")
                ax.axis("off")
                fig.savefig(output_path)
                plt.close(fig)
            with mock.patch("scripts.build_pptx.generate_summary_figures", return_value=None):
                build_pptx(config, logger)
            from pptx import Presentation

            prs = Presentation(str(paths["pptx_output_path"]))
            slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
            self.assertIn("summary_window_mode: configured_windows", slide_text)
            self.assertIn("baseline_window: -60 to 0 s", slide_text)
            self.assertIn("light_window: 5 to 20 s", slide_text)
            self.assertIn("post_window: 25 to 85 s", slide_text)
            self.assertNotIn("post_window: 15 to 75 s", slide_text)


if __name__ == "__main__":
    unittest.main()
